
import os
from pathlib import Path
from langchain_text_splitters import RecursiveCharacterTextSplitter
import chromadb
from chromadb.utils.embedding_functions import SentenceTransformerEmbeddingFunction
from config import settings

print("Initialisation ChromaDB...")
_embedding_fn = SentenceTransformerEmbeddingFunction(
    model_name=settings.EMBEDDING_MODEL
)
_client = chromadb.PersistentClient(path=settings.CHROMA_PATH)
print(f"ChromaDB pret")

_splitter = RecursiveCharacterTextSplitter(
    chunk_size=512, chunk_overlap=50,
    separators=["\n\n", "\n", ".", " "]
)

def _get_collection(course_id: str):
    # Chaque cours a sa propre collection isolee
    return _client.get_or_create_collection(
        name=f"cours_{course_id}",
        embedding_function=_embedding_fn,
        metadata={"hnsw:space": "cosine"}
    )

def _extract_pdf(path: str) -> str:
    import fitz
    doc = fitz.open(path)
    text = "\n".join(page.get_text() for page in doc)
    doc.close()
    return text

def _extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)

def _extract_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def _extract_txt(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")

def _extract_audio(path: str) -> str:
    from faster_whisper import WhisperModel
    print(f"  Transcription audio...")
    model = WhisperModel(settings.WHISPER_MODEL_SIZE, device="cpu", compute_type="int8")
    segments, info = model.transcribe(path, language=settings.WHISPER_LANGUAGE, beam_size=5)
    return " ".join(s.text for s in segments).strip()

def _extract_video(path: str) -> str:
    import subprocess, tempfile
    with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as f:
        tmp = f.name
    try:
        subprocess.run(
            ["ffmpeg", "-i", path, "-vn", "-acodec", "pcm_s16le",
             "-ar", "16000", "-ac", "1", tmp, "-y"],
            capture_output=True, timeout=300
        )
        return _extract_audio(tmp)
    finally:
        if os.path.exists(tmp): os.unlink(tmp)


def _extract_image(path: str) -> str:
    """OCR sur image (PNG, JPG, etc.)"""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(path)
        # Essayer français d abord, puis anglais
        try:
            text = pytesseract.image_to_string(img, lang="fra")
        except Exception:
            text = pytesseract.image_to_string(img)
        text = text.strip()
        if len(text) < 10:
            return f"Image : {path} (texte non extractible)"
        print(f"  OCR extrait : {len(text)} caracteres")
        return text
    except Exception as e:
        return f"Image non lisible : {e}"

_EXTRACTORS = {
    ".pdf": _extract_pdf, ".pptx": _extract_pptx, ".ppt": _extract_pptx,
    ".docx": _extract_docx, ".doc": _extract_docx,
    ".txt": _extract_txt, ".md": _extract_txt, ".csv": _extract_txt,
    ".mp3": _extract_audio, ".wav": _extract_audio, ".ogg": _extract_audio,
    ".mp4": _extract_video, ".mkv": _extract_video, ".avi": _extract_video, ".mov": _extract_video, ".webm": _extract_video,
    ".png": _extract_image, ".jpg": _extract_image, ".jpeg": _extract_image, ".webp": _extract_image, ".bmp": _extract_image,
}

def ingest_file_for_course(file_path: str, course_id: str) -> dict:
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix not in _EXTRACTORS:
        return {"status": "ignored", "reason": f"Format {suffix} non supporte", "chunks": 0}
    print(f"  Ingestion {path.name} pour cours {course_id}...")
    try:
        text = _EXTRACTORS[suffix](str(path))
    except Exception as e:
        return {"status": "error", "reason": str(e), "chunks": 0}
    if not text or not text.strip():
        return {"status": "error", "reason": "Document vide", "chunks": 0}
    chunks = _splitter.split_text(text)
    if not chunks:
        return {"status": "error", "reason": "Aucun chunk", "chunks": 0}
    col = _get_collection(course_id)
    ids = [f"{path.stem}_{i}" for i in range(len(chunks))]
    # Supprimer anciens chunks du meme fichier
    try:
        existing = col.get(where={"source": path.name})
        if existing["ids"]: col.delete(ids=existing["ids"])
    except Exception:
        pass
    col.add(
        documents=chunks,
        metadatas=[{"source": path.name, "type": suffix}] * len(chunks),
        ids=ids
    )
    print(f"  OK : {len(chunks)} chunks indexes")
    return {"status": "ok", "filename": path.name, "chunks": len(chunks)}

def search_in_course(query: str, course_id: str, k: int = 5) -> list[dict]:
    col = _get_collection(course_id)
    count = col.count()
    if count == 0:
        return []
    results = col.query(
        query_texts=[query],
        n_results=min(k, count)
    )
    if not results["documents"] or not results["documents"][0]:
        return []
    return [
        {"text": doc, "source": meta.get("source", ""), "distance": dist}
        for doc, meta, dist in zip(
            results["documents"][0],
            results["metadatas"][0],
            results["distances"][0]
        )
    ]

def get_context_for_course(query: str, course_id: str, k: int = 5) -> str:
    results = search_in_course(query, course_id, k)
    if not results:
        return ""
    return "\n\n---\n\n".join(r["text"] for r in results)

# Garder compatibilite avec l ancien code
def ingest_text(text: str, metadata: dict) -> int:
    col = _client.get_or_create_collection(
        name="formation_global",
        embedding_function=_embedding_fn
    )
    chunks = _splitter.split_text(text)
    if chunks:
        col.add(documents=chunks, metadatas=[metadata]*len(chunks),
                ids=[f"global_{i}" for i in range(len(chunks))])
    return len(chunks)

def ingest_file(file_path: str, metadata: dict = {}) -> dict:
    return ingest_file_for_course(file_path, "global")

def search(query: str, k: int = 5, source_filter: str = None) -> list[dict]:
    return search_in_course(query, "global", k)

def get_context_text(query: str, k: int = 5) -> str:
    return get_context_for_course(query, "global", k)

def list_sources(course_id: str = "global") -> list[dict]:
    try:
        col = _get_collection(course_id)
        all_items = col.get()
        sources = {}
        for meta in all_items.get("metadatas", []):
            src = meta.get("source", "inconnu")
            if src not in sources:
                sources[src] = {"source": src, "type": meta.get("type",""), "chunks": 0}
            sources[src]["chunks"] += 1
        return list(sources.values())
    except Exception:
        return []

def collection_count(course_id: str = "global") -> int:
    try:
        return _get_collection(course_id).count()
    except Exception:
        return 0
